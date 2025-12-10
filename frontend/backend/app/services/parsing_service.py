"""
Lightweight document parsing service for Vercel deployment.
Uses pypdf and python-docx instead of heavy ML libraries.
"""
import tempfile
import os
from pathlib import Path
from typing import Optional

# Lightweight libraries
from pypdf import PdfReader
from docx import Document as DocxDocument
import pptx

class LightweightParser:
    """
    Vercel-friendly parser.
    Avoids PyTorch/Docling to keep bundle size < 250MB.
    """
    
    def __init__(self):
        pass
    
    async def parse_document(self, file_content: bytes, filename: str, mime_type: str) -> str:
        """
        Parse document using lightweight libraries.
        """
        try:
            # Create temp file
            with tempfile.NamedTemporaryFile(delete=False, suffix=Path(filename).suffix) as tmp:
                tmp.write(file_content)
                tmp_path = tmp.name
            
            try:
                ext = Path(filename).suffix.lower()
                text = ""
                
                if ext == '.pdf':
                    reader = PdfReader(tmp_path)
                    for page in reader.pages:
                        text += page.extract_text() + "\n\n"
                
                elif ext in ['.docx', '.doc']:
                    doc = DocxDocument(tmp_path)
                    for para in doc.paragraphs:
                        text += para.text + "\n"
                        
                elif ext == '.pptx':
                    prs = pptx.Presentation(tmp_path)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                                
                elif ext in ['.txt', '.md']:
                    with open(tmp_path, 'r', encoding='utf-8', errors='ignore') as f:
                        text = f.read()
                        
                else:
                    text = f"[Unsupported file type: {ext}]\n"

                return text if text.strip() else "[Empty document]"
                
            finally:
                if os.path.exists(tmp_path):
                    os.unlink(tmp_path)
                    
        except Exception as e:
            print(f"Parsing failed: {e}")
            return f"[Error parsing file: {str(e)}]"

# Singleton instance
docling_parser = LightweightParser() # Keeping same variable name for compatibility
